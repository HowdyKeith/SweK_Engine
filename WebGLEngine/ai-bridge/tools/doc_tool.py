#!/usr/bin/env python3
# doc_tool.py - v1028  (document create/read for the doc bridge)
#   --op create --kind xlsx|docx|pptx|pdf  (--spec spec.json | --stdin)  --out FILE
#   --op read   --kind xlsx|docx|pptx|pdf  --path FILE  [--max N]
# Emits a single JSON line. xlsx is the priority (Mercor Excel work): create
# supports formulas (cells starting "="); read returns both formula + cached value.
import argparse, json, sys, os


def load_spec(a):
    if a.stdin:
        return json.loads(sys.stdin.read() or "{}")
    if a.spec:
        with open(a.spec, encoding="utf-8") as f:
            return json.load(f)
    return {}


# ---------------- CREATE ----------------
def create_xlsx(spec, out):
    from openpyxl import Workbook
    from openpyxl.styles import Font
    wb = Workbook(); wb.remove(wb.active)
    sheets = spec.get("sheets") or [{"name": "Sheet1", "rows": spec.get("rows", [])}]
    for sh in sheets:
        ws = wb.create_sheet(str(sh.get("name", "Sheet"))[:31])
        rows = sh.get("rows", [])
        for r in rows:
            ws.append(list(r))                      # str starting "=" -> formula automatically
        if sh.get("header") and rows:
            for c in ws[1]:
                c.font = Font(bold=True)
        for col, w in (sh.get("widths") or {}).items():
            ws.column_dimensions[str(col)].width = float(w)
    if not wb.sheetnames:
        wb.create_sheet("Sheet1")
    wb.save(out)
    return {"sheets": wb.sheetnames}


def create_docx(spec, out):
    from docx import Document
    doc = Document()
    if spec.get("title"):
        doc.add_heading(str(spec["title"]), level=0)
    blocks = spec.get("blocks")
    if blocks is None:
        blocks = [{"type": "para", "text": p} for p in spec.get("paragraphs", [])]
    for b in blocks:
        t = b.get("type", "para")
        if t == "heading":
            doc.add_heading(str(b.get("text", "")), level=int(b.get("level", 1)))
        elif t == "table":
            rows = b.get("rows", [])
            if rows:
                tbl = doc.add_table(rows=len(rows), cols=max(len(r) for r in rows))
                tbl.style = "Table Grid"
                for i, r in enumerate(rows):
                    for j, v in enumerate(r):
                        tbl.cell(i, j).text = str(v)
        else:
            p = doc.add_paragraph()
            run = p.add_run(str(b.get("text", "")))
            run.bold = bool(b.get("bold")); run.italic = bool(b.get("italic"))
    doc.save(out)
    return {"blocks": len(blocks)}


def create_pptx(spec, out):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    for sl in spec.get("slides", []):
        bullets = sl.get("bullets")
        layout = prs.slide_layouts[1] if bullets else prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = str(sl.get("title", ""))
        if bullets:
            body = slide.placeholders[1].text_frame
            for i, b in enumerate(bullets):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = str(b)
        elif sl.get("text"):
            tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.5), Inches(4)).text_frame
            tb.word_wrap = True; tb.text = str(sl["text"])
    prs.save(out)
    return {"slides": len(spec.get("slides", []))}


def create_pdf(spec, out):
    from fpdf import FPDF
    pdf = FPDF(); pdf.set_auto_page_break(True, 15); pdf.add_page()
    if spec.get("title"):
        pdf.set_font("Helvetica", "B", 18); pdf.multi_cell(0, 10, str(spec["title"])); pdf.ln(2)
    blocks = spec.get("blocks")
    if blocks is None:
        blocks = [{"type": "para", "text": p} for p in spec.get("paragraphs", [])]
    for b in blocks:
        txt = str(b.get("text", "")).encode("latin-1", "replace").decode("latin-1")
        if b.get("type") == "heading":
            pdf.set_font("Helvetica", "B", 14)
        else:
            pdf.set_font("Helvetica", "", 11)
        pdf.multi_cell(0, 7, txt); pdf.ln(1)
    pdf.output(out)
    return {"blocks": len(blocks)}


# ---------------- READ ----------------
def read_xlsx(path, mx):
    from openpyxl import load_workbook
    wbf = load_workbook(path, data_only=False)
    try:
        wbv = load_workbook(path, data_only=True)
    except Exception:
        wbv = None
    sheets = []
    for name in wbf.sheetnames[:8]:
        wsf = wbf[name]; wsv = wbv[name] if wbv else None
        cells = []
        for row in wsf.iter_rows():
            for c in row:
                if c.value is None:
                    continue
                cell = {"ref": c.coordinate, "value": c.value}
                if isinstance(c.value, str) and c.value.startswith("="):
                    cell["formula"] = c.value
                    if wsv is not None:
                        cv = wsv[c.coordinate].value
                        cell["cached"] = cv
                cells.append(cell)
                if len(cells) >= mx:
                    break
            if len(cells) >= mx:
                break
        sheets.append({"name": name, "maxRow": wsf.max_row, "maxCol": wsf.max_column, "cellCount": len(cells), "cells": cells})
    return {"sheets": sheets, "note": "cached = last value Excel saved; openpyxl does not recompute formulas."}


def read_docx(path, mx):
    from docx import Document
    doc = Document(path)
    paras = [{"text": p.text, "style": p.style.name if p.style else ""} for p in doc.paragraphs if p.text.strip()][:mx]
    tables = []
    for t in doc.tables[:20]:
        tables.append([[c.text for c in r.cells] for r in t.rows])
    return {"paragraphs": paras, "tables": tables}


def read_pptx(path, mx):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, sl in enumerate(prs.slides):
        texts = []
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts.append(sh.text_frame.text)
        slides.append({"slide": i + 1, "texts": texts})
        if i + 1 >= mx:
            break
    return {"slideCount": len(slides), "slides": slides}


def read_pdf(path, mx):
    from pypdf import PdfReader
    r = PdfReader(path)
    pages = []
    for i, pg in enumerate(r.pages):
        if i >= mx:
            break
        try:
            pages.append({"page": i + 1, "text": (pg.extract_text() or "")[:4000]})
        except Exception as e:
            pages.append({"page": i + 1, "error": str(e)})
    return {"pageCount": len(r.pages), "pages": pages}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--op", required=True, choices=["create", "read"])
    ap.add_argument("--kind", required=True, choices=["xlsx", "docx", "pptx", "pdf"])
    ap.add_argument("--spec"); ap.add_argument("--stdin", action="store_true")
    ap.add_argument("--path"); ap.add_argument("--out"); ap.add_argument("--max", type=int, default=300)
    a = ap.parse_args()
    try:
        if a.op == "create":
            if not a.out:
                print(json.dumps({"ok": False, "error": "--out required for create"})); return
            spec = load_spec(a)
            fn = {"xlsx": create_xlsx, "docx": create_docx, "pptx": create_pptx, "pdf": create_pdf}[a.kind]
            info = fn(spec, a.out)
            print(json.dumps({"ok": True, "op": "create", "kind": a.kind, "out": a.out, "bytes": os.path.getsize(a.out), **info}))
        else:
            if not a.path or not os.path.exists(a.path):
                print(json.dumps({"ok": False, "error": "file not found: " + str(a.path)})); return
            fn = {"xlsx": read_xlsx, "docx": read_docx, "pptx": read_pptx, "pdf": read_pdf}[a.kind]
            info = fn(a.path, a.max)
            print(json.dumps({"ok": True, "op": "read", "kind": a.kind, "path": a.path, **info}, default=str))
    except Exception as e:
        import traceback
        print(json.dumps({"ok": False, "op": a.op, "kind": a.kind, "error": str(e), "trace": traceback.format_exc()[-500:]}))


if __name__ == "__main__":
    main()
